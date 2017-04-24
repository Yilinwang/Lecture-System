import os
import django
import subprocess
from collections import defaultdict
from datetime import timedelta

os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'lsite.settings')
django.setup()
from lecture.models import Slide
from lecture.models import Slidekeyterm
from lecture.models import KeytermRelation
from lecture.models import VideoAttr
from lecture.models import ChapVideoAttr
from lecture.models import SumAttr
from lecture.models import SumPageTitle

from glob import glob
import pptx

def clean():
    #Slide.objects.all().delete()
    KeytermRelation.objects.all().delete()
    #Slidekeyterm.objects.all().delete()

def init(num):
    for x in glob('./lecture/static/lecture/slides/'+num+'/*'):
        a = x.split(num+'/')[1][:-4]
        b = a.split('-')
        print(b)
        tmp = Slide(chapter=int(num), subchapter=int(b[0]), page=int(b[1]))
        tmp.save()

def addkeyterm(t, k):
    print(t, k)
    for x in k:
        tmp = Slidekeyterm(title=t, keyterm=x)
        tmp.save()

def addrelation(a, b):
    tmp = KeytermRelation(k1=a, k2=b)
    tmp.save()
    #tmp = KeytermRelation(k1=b, k2=a)
    #tmp.save()

def percentage(a):
    if a == 0:
        return 0
    p = round(a*100)
    return p if p > 0 else 1

def sumvideo():
    s = defaultdict(lambda: defaultdict(list))
    bs = defaultdict(lambda: defaultdict(list))
    for x in glob('lecture/static/lecture/videos/summary/*.mp4'):
        title = x.split('/')[-1].split('.')[0]
        time = subprocess.run(['ffprobe', x], stderr=subprocess.PIPE).stderr.strip().split(b'Duration: ')[1].split(b',')[0].split(b'.')[0].split(b':', 1)[1].decode('ascii')
        s[title] = time
    for x in glob('lecture/static/lecture/videos/brief_summary/*.mp4'):
        title = x.split('/')[-1].split('.')[0]
        time = subprocess.run(['ffprobe', x], stderr=subprocess.PIPE).stderr.strip().split(b'Duration: ')[1].split(b',')[0].split(b'.')[0].split(b':', 1)[1].decode('ascii')
        bs[title] = time
    for t in s:
        tmp = SumAttr(title=t, time=s[t], brief_time=bs[t])
        print(tmp.title, tmp.time, tmp.brief_time)
        tmp.save()

def chapvideo():
    print('hi')
    prev = './lecture/static/lecture/videos/'
    for x in glob('./lecture/static/lecture/videos/chapter3/*'):
        title = x.strip().split('/')[-1]
        title_text = title[:-4]
        print(title_text)
        chapter3 = subprocess.run(['ffprobe', x], stderr=subprocess.PIPE).stderr.strip().split(b'Duration: ')[1].split(b',')[0].split(b'.')[0].decode('ascii')
        chapter10 = subprocess.run(['ffprobe', prev+'/chapter10/'+title], stderr=subprocess.PIPE).stderr.strip().split(b'Duration: ')[1].split(b',')[0].split(b'.')[0].decode('ascii')
        time = subprocess.run(['ffprobe', prev+'/whole/'+title], stderr=subprocess.PIPE).stderr.strip().split(b'Duration: ')[1].split(b',')[0].split(b'.')[0].decode('ascii')
        tmp = ChapVideoAttr(title=title_text, time=time, chapter3time=chapter3, chapter10time=chapter10)
        tmp.save()

def video():
    d = defaultdict(lambda: defaultdict(list))
    for x in glob('./lecture/static/lecture/videos/*'):
        title = x.split('/')[-1].split('.')[0]
        time = subprocess.run(['ffprobe', x], stderr=subprocess.PIPE).stderr.strip().split(b'Duration: ')[1].split(b',')[0].split(b'.')[0].split(b':', 1)[1].decode('ascii')
        sec = int(time.split(':')[0])*60+int(time.split(':')[1])
        d[int(title.split('-')[0])][int(title.split('-')[1])].append((title, time, sec))

    allsum = 0
    chsum = defaultdict(int)
    subchsum = defaultdict(lambda: defaultdict(int))
    for x in d:
        for y in d[x]:
            d[x][y] = sorted(d[x][y], key=lambda k: int(k[0].split('-')[2]))
            s = sum([k[2] for k in d[x][y]])
            chsum[x] += s
            subchsum[x][y] += s
            allsum += s
    prev_ch = 0
    for x in d:
        prev_subch = 0
        for y in d[x]:
            prev = 0
            for z in d[x][y]:
                print(z[0], percentage(prev_ch/allsum), percentage(chsum[x]/allsum), str(timedelta(seconds=chsum[x])), percentage(prev_subch/chsum[x]), percentage(subchsum[x][y]/chsum[x]), str(timedelta(seconds=subchsum[x][y])), percentage(prev/subchsum[x][y]), percentage(z[2]/subchsum[x][y]), z[1])
                tmp = VideoAttr(title=z[0], prev_ch=percentage(prev_ch/allsum), cur_ch=percentage(chsum[x]/allsum), ch_time=str(timedelta(seconds=chsum[x])), prev_subch=percentage(prev_subch/chsum[x]), cur_subch=percentage(subchsum[x][y]/chsum[x]), subch_time=str(timedelta(seconds=subchsum[x][y])), prev=percentage(prev/subchsum[x][y]), cur=percentage(z[2]/subchsum[x][y]), time=z[1])
                tmp.save()
                prev += z[2]
            prev_subch += subchsum[x][y]
        prev_ch += chsum[x]

def m():
    tmp = VideoAttr(title='1', prev_ch=1, cur_ch=1, ch_time='1', prev_subch=1, cur_subch=1, subch_time='1', prev = 1, cur = 1, time='1')

def text_title():
    l = sorted(list(Slide.objects.filter(chapter=10)), key=lambda x: (x.chapter, x.subchapter, x.page))
    p = pptx.Presentation('/Users/yilin/Documents/NTU/16fall/lecture_sys/material/Slide/ppt/course10.0.pptx')
    i = 0
    for s in p.slides:
        title = ''
        for x in s.shapes:
            if type(x) == pptx.shapes.autoshape.Shape:
                title = x.text.split('\n')[0].strip()
                break
        print(title, l[i].chapter, l[i].subchapter, l[i].page)
        l[i].title = title
        l[i].save()
        i += 1
    print(len(l), len(p.slides))

def import_keyterm():
    with open('keyterm.txt') as fp:
        for line in fp:
            k = line.strip().split(' ')
            if len(k) > 1:
                tmp = k[0].split('-')
                addkeyterm(tmp[0], k[1:])
                addkeyterm('%s-%s' % (tmp[0], tmp[1]), k[1:])
                addkeyterm(k[0], k[1:])

def import_keyterm_graph():
    with open('keyterm_graph.txt') as fp:
        cur = ''
        for line in fp:
            if '\t' not in line and line != '\n':
                cur = line.strip()
            elif '\t' in line:
                print(cur, line.strip())
                addrelation(cur, line.strip())

def saveSumPageTitle(t, tt):
    tmp = SumPageTitle(title=t, title_text=tt)
    tmp.save()

def addSumPageTitle():
    saveSumPageTitle('2', 'Linear Time-invariant Systems')
    saveSumPageTitle('2-1', 'Discrete-time Systems: the Convolution Sum')
    saveSumPageTitle('2-2', 'Vector Space Interpretation for Discrete-time Systems')
    saveSumPageTitle('2-3', 'Continuous-time System : the Convolution Integral')
    saveSumPageTitle('2-4', 'Properties of Linear Time-invariant Systems')
    saveSumPageTitle('2-5', 'Systems Described by Differential/Difference Equations')
    saveSumPageTitle('2-6', 'The Unit Impulse for Continuous-time Cases')
    saveSumPageTitle('2-7', 'Vector Space Interpretation for Continuous-time Systems')
    saveSumPageTitle('3', 'Fourier Series Representation of Periodic Signals')
    saveSumPageTitle('3-1', 'Exponential/Sinusoidal Signals as Building Blocks for Many Signals')
    saveSumPageTitle('3-2', 'Fourier Series Representation of Continuous-time Periodic Signals')
    saveSumPageTitle('3-3', 'Properties of Fourier Series')
    saveSumPageTitle('3-4', 'Fourier Series Representation of Discrete-time Periodic Signals')
    saveSumPageTitle('3-5', 'Application Example')
    saveSumPageTitle('4', 'Continuous-time Fourier Transform')
    saveSumPageTitle('4-1', 'From Fourier Series to Fourier Transform')
    saveSumPageTitle('4-2', 'Properties of Continuous-time Fourier Transform')
    saveSumPageTitle('5', 'Discrete-time Fourier Transform')
    saveSumPageTitle('5-1', 'Discrete-time Fourier Transform Representation for discrete-time signals')
    saveSumPageTitle('5-2', 'Properties of Discrete-time Fourier Transform')
    saveSumPageTitle('5-3', 'Summary and Duality')
    saveSumPageTitle('6', 'Time/Frequency Characterization of Signals/Systems')
    saveSumPageTitle('6-1', 'Magnitude and Phase for Signals and Systems')
    saveSumPageTitle('6-2', 'Filtering')
    saveSumPageTitle('6-3', 'First/Second-Order Systems Described by Differential/Difference Equations')
    saveSumPageTitle('7', 'Sampling')
    saveSumPageTitle('7-1', 'The Sampling Theorem')
    saveSumPageTitle('7-2', 'Discrete-time Processing of Continuous-time Signals')
    saveSumPageTitle('7-3', 'Change of Sampling Frequency')
    saveSumPageTitle('8', 'Communication Systems')
    saveSumPageTitle('8-0', 'Communication Systems')
    saveSumPageTitle('8-1', 'Amplitude Modulation (AM) and Frequency-Division Multiplexing (FDM)')
    saveSumPageTitle('8-2', 'Pulse Modulation and Time-Division Multiplexing')
    saveSumPageTitle('8-3', 'Angle/Frequency Modulation')
    saveSumPageTitle('8-4', 'Discrete-time Modulation')
    saveSumPageTitle('9', 'Laplace Transform')
    saveSumPageTitle('9-1', 'General Principles of Laplace Transform')
    saveSumPageTitle('9-2', 'Properties of Laplace Transform')
    saveSumPageTitle('9-3', 'System Characterization with Laplace Transform')
    saveSumPageTitle('9-4', 'Unilateral Laplace Transform')
    saveSumPageTitle('10', 'Z-Transform')
    saveSumPageTitle('10-1', 'General Principles of Z-Transform')
    saveSumPageTitle('10-2', 'Properties of Z-Transform')
    saveSumPageTitle('10-3', 'System Characterization with Z-Transform')
    saveSumPageTitle('10-4', 'Unilateral Z-Transform')

def main():
    #Slidekeyterm.objects.all().delete()
    #import_keyterm()
    #chapvideo()
    pass

if __name__ == '__main__':
    main()
